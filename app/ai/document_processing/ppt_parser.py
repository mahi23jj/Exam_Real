import io
from pptx import Presentation
from typing import List, Dict, Any


class PPTParser:
    """Extracts text per slide from PPT/PPTX presentations using python-pptx."""

    @staticmethod
    def parse_ppt_bytes(ppt_bytes: bytes) -> List[Dict[str, Any]]:
        """Parses PPT/PPTX bytes and returns list of slide dicts: [{'page_number': 1, 'text': '...'}]"""
        slides_content = []
        prs = Presentation(io.BytesIO(ppt_bytes))

        for slide_index, slide in enumerate(prs.slides):
            slide_text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text_parts.append(shape.text.strip())

            full_slide_text = "\n".join(slide_text_parts)
            slides_content.append({
                "page_number": slide_index + 1,
                "text": full_slide_text
            })

        return slides_content
